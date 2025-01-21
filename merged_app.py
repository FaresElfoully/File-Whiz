import os
import shutil
import pandas as pd
from pptx import Presentation
from docx import Document
from langchain_community.vectorstores import FAISS
from langchain.prompts import ChatPromptTemplate
from groq import Groq
from langchain_community.document_loaders import DirectoryLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.schema.document import Document
from langchain_ollama import OllamaEmbeddings
import faiss
import warnings
from PyPDF2 import PdfReader
import pickle
import numpy as np
import datetime
import cohere
from rank_bm25 import BM25Okapi
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
import nltk
from concurrent.futures import ThreadPoolExecutor
import threading
from dotenv import load_dotenv
import re
import traceback

# Load environment variables
load_dotenv()

# Constants
FAISS_PATH = "faiss_index"
DATA_PATH = "outputs"
CHUNKS_PATH = os.path.join(FAISS_PATH, "chunks.pkl")
COHERE_API_KEY = os.getenv('COHERE_API_KEY', '3hLheqT2lFX5RjiSEOzsu04rKStDF1LSyMwTzTIk')
HNSW_M = 16  # Number of connections per layer
HNSW_EF_CONSTRUCTION = 200  # Size of dynamic candidate list for construction
HYBRID_ALPHA = 0.7  # Weight for dense embeddings vs sparse search

def detect_language(text):
    """Detect if text contains Arabic script."""
    arabic_pattern = re.compile('[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF\uFB50-\uFDFF\uFE70-\uFEFF]')
    return 'arabic' if arabic_pattern.search(text) else 'english'

PROMPT_TEMPLATE = """
You are a helpful assistant that provides accurate and relevant information. 
IMPORTANT: If the input question is in Arabic, you MUST respond in Arabic. If in English, respond in English.

Answer the question based only on the following context:

{context}

---

Question: {question}

Response Language Instructions:
1. If the question contains Arabic text -> RESPOND IN ARABIC ONLY
2. If the question is in English -> Respond in English
3. Match the exact language of the question

Important Instructions:
1. Answer ONLY based on the provided context
2. Maintain the same tone and style as the question
3. Be concise and direct in your response
"""

EXCEL_PROMPT_TEMPLATE = """
You are an expert data analyst specializing in Excel data analysis. 
Analyze and answer questions about the following Excel data:

Context:
{context}

Question: {question}

Important Instructions:
1. If the question is in Arabic, provide the complete analysis in Arabic
2. If the question is in English, provide the analysis in English
3. Consider:
   - Data patterns and trends
   - Statistical insights
   - Data quality
   - Relationships between columns
   - Numerical calculations
   - Business implications
4. Structure your answer to:
   - Directly address the question
   - Support conclusions with data
   - Highlight key patterns
   - Provide actionable insights
   - Note any limitations
"""

# Add this function after the imports and constants, before the other functions
def preprocess_text(text: str) -> list[str]:
    """Preprocess text with support for Arabic and other languages."""
    try:
        nltk.data.find('tokenizers/punkt')
        nltk.data.find('corpora/stopwords')
    except LookupError:
        nltk.download('punkt')
        nltk.download('stopwords')
    
    # Handle None or empty text
    if not text:
        return []
    
    # Convert to string if needed
    if not isinstance(text, str):
        text = str(text)
    
    # Basic tokenization that works with Arabic and other scripts
    tokens = word_tokenize(text)
    
    # Remove punctuation and normalize
    tokens = [token.lower() for token in tokens if token.isalnum()]
    
    # For Arabic text, we might want to skip stemming
    if any('\u0600' <= char <= '\u06FF' for char in text):  # Arabic Unicode range
        return tokens
    
    # For English text, proceed with stopword removal and stemming
    stop_words = set(stopwords.words('english'))
    stemmer = PorterStemmer()
    tokens = [
        stemmer.stem(token)
        for token in tokens
        if token not in stop_words and token.isalpha()
    ]
    
    return tokens

# Embedding function
def get_embedding_function():
    embeddings = OllamaEmbeddings(model="nomic-embed-text")
    return embeddings

# File conversion functions
def get_user_input(prompt_message: str, default_value: str = None) -> str:
    if default_value:
        user_input = input(f"{prompt_message} (press Enter for default: {default_value}): ").strip()
        return user_input if user_input else default_value
    return input(f"{prompt_message}: ").strip()

def convert_files_to_md(input_dir=None, output_dir=None):
    print("\n🔍 Let's convert some files to markdown!")
    
    if input_dir is None:
        input_dir = get_user_input("📁 Enter the input directory path containing your files")
    
    if output_dir is None:
        output_dir = get_user_input("📂 Enter the output directory path for markdown files", "outputs")
    
    if not os.path.exists(input_dir):
        print("❌ Input directory doesn't exist!")
        return
    
    os.makedirs(output_dir, exist_ok=True)
    print(f"\n🚀 Starting conversion from {input_dir} to {output_dir}...")
    print(f"Found files in input directory: {os.listdir(input_dir)}")

    files_processed = 0
    for filename in os.listdir(input_dir):
        input_path = os.path.join(input_dir, filename)
        print(f"\n📄 Processing file: {filename}")
        
        if filename.endswith(".pdf"):
            print("Converting PDF to markdown...")
            convert_pdf_to_md(input_path, output_dir)
            files_processed += 1
        elif filename.endswith(".xlsx"):
            print("Converting Excel to markdown...")
            convert_excel_to_md(input_path, output_dir)
            files_processed += 1
        elif filename.endswith(".pptx"):
            print("Converting PowerPoint to markdown...")
            convert_pptx_to_md(input_path, output_dir)
            files_processed += 1
        elif filename.endswith(".docx"):
            print("Converting Word to markdown...")
            convert_word_to_md(input_path, output_dir)
            files_processed += 1
        else:
            print(f"⚠️ Unsupported file type: {filename}")
    
    print(f"\n✅ Conversion complete! Processed {files_processed} files")
    if files_processed == 0:
        print("⚠️ Warning: No files were converted to markdown")

def convert_pdf_to_md(input_path, output_dir):
    reader = PdfReader(input_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    save_md_file(output_dir, os.path.basename(input_path), text)

def convert_excel_to_md(input_path, output_dir):
    """Enhanced Excel to Markdown conversion with comprehensive data analysis and formatting."""
    try:
        print(f"\n📊 Analyzing Excel file: {input_path}")
        
        # Read all sheets with enhanced options
        print("Reading Excel file...")
        excel_data = pd.read_excel(
            input_path, 
            sheet_name=None,
            na_filter=True,
            dtype=str,  # Read all as string to preserve formatting
            engine='openpyxl'
        )
        
        text = []
        
        # File Information Section
        text.append("# 📊 Comprehensive Excel File Analysis Report")
        
        # File Metadata
        file_stats = os.stat(input_path)
        created_time = datetime.datetime.fromtimestamp(file_stats.st_ctime)
        modified_time = datetime.datetime.fromtimestamp(file_stats.st_mtime)
        
        text.append("\n## 📌 File Metadata")
        text.append(f"- 📄 Filename: `{os.path.basename(input_path)}`")
        text.append(f"- 📁 Full Path: `{os.path.abspath(input_path)}`")
        text.append(f"- 📅 Created: {created_time.strftime('%Y-%m-%d %H:%M:%S')}")
        text.append(f"- 🔄 Last Modified: {modified_time.strftime('%Y-%m-%d %H:%M:%S')}")
        text.append(f"- 💾 File Size: {file_stats.st_size / 1024:.2f} KB")
        text.append(f"- 📑 Number of Sheets: {len(excel_data)}")
        
        # Overall File Statistics
        total_rows = sum(df.shape[0] for df in excel_data.values())
        total_cols = sum(df.shape[1] for df in excel_data.values())
        total_cells = sum(df.size for df in excel_data.values())
        
        text.append("\n## 📊 Overall File Statistics")
        text.append(f"- 📈 Total Rows (all sheets): {total_rows:,}")
        text.append(f"- 📉 Total Columns (all sheets): {total_cols:,}")
        text.append(f"- 🔢 Total Cells: {total_cells:,}")
        
        # Process each sheet
        for sheet_name, df in excel_data.items():
            print(f"\nProcessing sheet: {sheet_name}")
            
            text.append(f"\n## 📋 Sheet Analysis: {sheet_name}")
            
            # Sheet Overview
            rows, cols = df.shape
            text.append("\n### 📊 Sheet Overview")
            text.append(f"- 📏 Dimensions: {rows:,} rows × {cols} columns")
            text.append(f"- 🔢 Total Cells: {rows * cols:,}")
            
            # Column Analysis
            text.append("\n### 🔍 Column Analysis")
            for col in df.columns:
                text.append(f"\n#### 📊 Column: {col}")
                
                # Basic column statistics
                non_empty = df[col].count()
                empty = df[col].isna().sum()
                unique_values = df[col].nunique()
                
                text.append(f"- 📝 Non-empty Values: {non_empty:,} ({(non_empty/len(df)*100):.1f}%)")
                text.append(f"- ⚪ Empty Values: {empty:,} ({(empty/len(df)*100):.1f}%)")
                text.append(f"- 🎯 Unique Values: {unique_values:,} ({(unique_values/len(df)*100):.1f}% of total)")
                
                # Value frequency analysis
                if unique_values <= 20:  # Show distribution for categorical-like columns
                    text.append("\n📊 Value Distribution:")
                    value_counts = df[col].value_counts().head(10)
                    for val, count in value_counts.items():
                        text.append(f"  - {val}: {count:,} ({(count/len(df)*100):.1f}%)")
                
                # Text length analysis for string columns
                try:
                    length_stats = df[col].str.len().describe()
                    text.append("\n📏 Text Length Statistics:")
                    text.append(f"- Average Length: {length_stats['mean']:.1f} characters")
                    text.append(f"- Min Length: {length_stats['min']:.0f} characters")
                    text.append(f"- Max Length: {length_stats['max']:.0f} characters")
                except:
                    pass
                
                # Pattern Analysis
                try:
                    numeric_values = df[col].str.contains(r'\d').sum()
                    alpha_values = df[col].str.contains(r'[a-zA-Z]').sum()
                    special_chars = df[col].str.contains(r'[^a-zA-Z0-9\s]').sum()
                    
                    text.append("\n🔍 Content Patterns:")
                    text.append(f"- Contains Numbers: {numeric_values:,} values ({(numeric_values/len(df)*100):.1f}%)")
                    text.append(f"- Contains Letters: {alpha_values:,} values ({(alpha_values/len(df)*100):.1f}%)")
                    text.append(f"- Contains Special Characters: {special_chars:,} values ({(special_chars/len(df)*100):.1f}%)")
                except:
                    pass
            
            # Data Quality Assessment
            text.append("\n### ✅ Data Quality Assessment")
            quality_score = 100
            quality_issues = []
            
            # Check for duplicates
            duplicates = df.duplicated().sum()
            if duplicates > 0:
                quality_issues.append(f"- ⚠️ Found {duplicates:,} duplicate rows ({(duplicates/len(df)*100):.1f}%)")
                quality_score -= min(20, (duplicates/len(df)*100))
            
            # Check for missing values
            missing_cols = df.columns[df.isna().any()].tolist()
            if missing_cols:
                quality_issues.append(f"- ⚠️ {len(missing_cols)} columns have missing values:")
                for col in missing_cols:
                    missing_count = df[col].isna().sum()
                    quality_issues.append(f"  - {col}: {missing_count:,} missing ({(missing_count/len(df)*100):.1f}%)")
                quality_score -= min(30, len(missing_cols) * 5)
            
            text.append(f"\n📊 Data Quality Score: {max(0, quality_score):.1f}/100")
            if quality_issues:
                text.append("\n⚠️ Quality Issues:")
                text.extend(quality_issues)
            else:
                text.append("\n✅ No major quality issues detected")
            
            # Data Preview
            text.append("\n### 👀 Data Preview")
            
            # First rows
            text.append("\n#### ⬆️ First 5 Rows:")
            text.append(df.head().to_markdown(index=False, tablefmt="pipe"))
            
            # Last rows
            if len(df) > 5:
                text.append("\n#### ⬇️ Last 5 Rows:")
                text.append(df.tail().to_markdown(index=False, tablefmt="pipe"))
            
            # Complete Data
            text.append("\n### 📋 Complete Dataset")
            text.append("\n<details>")
            text.append("<summary>Click to expand full table</summary>\n")
            
            # Convert to markdown table with proper formatting
            try:
                markdown_table = df.to_markdown(
                    index=False,
                    tablefmt="pipe",
                    stralign="left",
                    numalign="right"
                )
                text.append(markdown_table)
            except Exception as e:
                print(f"Warning: Error converting full table to markdown: {str(e)}")
                # Fallback: Try converting smaller chunks
                chunk_size = 1000
                for i in range(0, len(df), chunk_size):
                    chunk = df.iloc[i:i + chunk_size]
                    try:
                        chunk_table = chunk.to_markdown(
                            index=False,
                            tablefmt="pipe",
                            stralign="left",
                            numalign="right"
                        )
                        text.append(chunk_table)
                        text.append("\n")
                    except Exception as e2:
                        print(f"Warning: Error converting table chunk: {str(e2)}")
            
            text.append("\n</details>")
            
            # Add separator between sheets
            text.append("\n---\n")
        
        # Save the markdown file
        final_text = "\n".join(text)
        save_md_file(output_dir, os.path.basename(input_path), final_text)
        print(f"✨ Enhanced Excel analysis completed for {input_path}")
        
    except Exception as e:
        print(f"Error converting Excel file {input_path}: {str(e)}")
        # Fallback method
        try:
            print("Attempting basic conversion method...")
            excel_data = pd.read_excel(input_path, sheet_name=None, dtype=str)
            text = ["# Excel File Content\n"]
            for sheet, data in excel_data.items():
                text.append(f"\n## Sheet: {sheet}\n")
                text.append(data.to_markdown(index=False) + "\n\n")
            save_md_file(output_dir, os.path.basename(input_path), "\n".join(text))
            print("✓ Basic conversion completed successfully")
        except Exception as e2:
            print(f"Fallback conversion also failed: {str(e2)}")

def convert_pptx_to_md(input_path, output_dir):
    presentation = Presentation(input_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    save_md_file(output_dir, os.path.basename(input_path), text)

def convert_word_to_md(input_path, output_dir):
    doc = Document(input_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    save_md_file(output_dir, os.path.basename(input_path), text)

def save_md_file(output_dir, original_filename, content):
    """Save markdown file with UTF-8 encoding."""
    base_name = os.path.splitext(original_filename)[0]
    output_path = os.path.join(output_dir, f"{base_name}.md")
    
    # Ensure content is properly encoded
    if isinstance(content, bytes):
        content = content.decode('utf-8', errors='ignore')
    
    try:
        with open(output_path, "w", encoding="utf-8") as md_file:
            md_file.write(content)
        print(f"✓ Converted: {original_filename} -> {output_path}")
    except Exception as e:
        print(f"Warning: Error saving file {original_filename}: {str(e)}")

# Database population functions
def load_documents():
    """Load documents with support for non-ASCII filenames and content."""
    documents = []
    
    for filename in os.listdir(DATA_PATH):
        try:
            file_path = os.path.join(DATA_PATH, filename)
            if not filename.endswith('.md'):
                continue
                
            # Read the markdown file with UTF-8 encoding
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                
            # Create document with proper metadata
            doc = Document(
                page_content=content,
                metadata={
                    "source": file_path,
                    "filename": filename
                }
            )
            documents.append(doc)
            
        except Exception as e:
            print(f"Warning: Error loading file {filename}: {str(e)}")
            continue
    
    if not documents:
        raise ValueError("No documents were successfully loaded")
        
    return documents

def split_documents(documents: list[Document]):
    """Enhanced document splitting with adaptive chunk sizes and categorization."""
    # Initialize categorizer
    categorizer = cohere.Client(COHERE_API_KEY)
    
    def get_optimal_chunk_size(text: str) -> int:
        """Determine optimal chunk size based on content complexity."""
        avg_word_length = sum(len(word) for word in text.split()) / len(text.split())
        if avg_word_length > 8:  # Complex technical content
            return 600
        elif avg_word_length < 4:  # Simple content
            return 1000
        return 800
    
    chunks = []
    for doc in documents:
        chunk_size = get_optimal_chunk_size(doc.page_content)
        
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_size // 10,
            length_function=len,
            is_separator_regex=False,
        )
        
        doc_chunks = text_splitter.split_documents([doc])
        
        # Categorize chunks in batches, silently handle errors
        for chunk in doc_chunks:
            try:
                response = categorizer.classify(
                    text=chunk.page_content,
                    examples=[
                        {"text": "Revenue grew by 50%", "label": "financial"},
                        {"text": "The protein structure", "label": "scientific"},
                        {"text": "The team decided", "label": "organizational"}
                    ],
                    model='embed-english-v3.0'
                )
                chunk.metadata['category'] = response.classifications[0].prediction
            except:
                # Silently set category to unknown without printing warnings
                chunk.metadata['category'] = 'unknown'
        
        chunks.extend(doc_chunks)
    
    return chunks

def calculate_chunk_ids(chunks):
    last_page_id = None
    current_chunk_index = 0

    for chunk in chunks:
        source = chunk.metadata.get("source")
        page = chunk.metadata.get("page")
        current_page_id = f"{source}:{page}"

        if current_page_id == last_page_id:
            current_chunk_index += 1
        else:
            current_chunk_index = 0

        chunk_id = f"{current_page_id}:{current_chunk_index}"
        last_page_id = current_page_id
        chunk.metadata["id"] = chunk_id

    return chunks

def add_to_faiss(chunks: list[Document]):
    """Add documents to FAISS using HNSW index with hybrid search capability."""
    embeddings = get_embedding_function()
    
    # Get dimension from first embedding
    first_embedding = embeddings.embed_query(chunks[0].page_content)
    dimension = len(first_embedding)
    
    # Initialize HNSW index
    index = faiss.IndexHNSWFlat(dimension, HNSW_M)
    index.hnsw.efConstruction = HNSW_EF_CONSTRUCTION
    
    # Prepare for hybrid search
    tokenized_chunks = [preprocess_text(chunk.page_content) for chunk in chunks]
    bm25 = BM25Okapi(tokenized_chunks)
    
    # Process embeddings in parallel
    with ThreadPoolExecutor() as executor:
        embeddings_list = list(executor.map(
            lambda chunk: embeddings.embed_query(chunk.page_content),
            chunks
        ))
    
    # Convert to numpy array and add to index
    embeddings_array = np.array(embeddings_list, dtype='float32')
    index.add(embeddings_array)
    
    # Save everything
    if not os.path.exists(FAISS_PATH):
        os.makedirs(FAISS_PATH)
    
    faiss.write_index(index, os.path.join(FAISS_PATH, "faiss_index.index"))
    
    # Save chunks and BM25 index
    with open(CHUNKS_PATH, 'wb') as f:
        pickle.dump({
            'chunks': chunks,
            'bm25': bm25,
            'tokenized_chunks': tokenized_chunks
        }, f)
    
    print(f"✨ Enhanced FAISS index saved with HNSW and hybrid search capability")

def clear_database():
    """Clears the FAISS database."""
    if os.path.exists(FAISS_PATH):
        print(f"⚠️  Deleting existing FAISS index at {FAISS_PATH}...")
        faiss_index_path = os.path.join(FAISS_PATH, "faiss_index.index")
        if os.path.exists(faiss_index_path):
            os.remove(faiss_index_path)
            print("✅ FAISS index deleted.")
        else:
            print("❌ FAISS index file not found.")
    else:
        print("❌ FAISS database directory not found.")

# Query functions
def query_faiss(query_text: str = None):
    if query_text is None:
        return {
            'error': 'No query text provided'
        }
    
    try:
        # Check if CHUNKS_PATH exists
        if not os.path.exists(CHUNKS_PATH):
            return {
                'error': 'No documents have been processed yet. Please upload and process documents first.'
            }
            
        # Load saved data
        with open(CHUNKS_PATH, 'rb') as f:
            saved_data = pickle.load(f)
            chunks = saved_data['chunks'] if isinstance(saved_data, dict) else saved_data
            bm25 = saved_data.get('bm25') if isinstance(saved_data, dict) else None
        
        # If no chunks were loaded
        if not chunks:
            return {
                'error': 'No document chunks found. Please process documents first.'
            }
            
        # Preprocess query
        processed_query = preprocess_text(query_text)
        
        # If we have only one chunk, use it directly
        if len(chunks) == 1:
            context_text = chunks[0].page_content
            source = chunks[0].metadata.get('source', 'combined_docs')
            sources = [{
                'name': os.path.basename(source),
                'location': 'Full Document',
                'content': context_text[:200] + "..." if len(context_text) > 200 else context_text,
                'relevance': 'Primary Source'
            }]
            
            return {
                'response': context_text,
                'sources': sources
            }
        
        # Hybrid search
        embeddings = get_embedding_function()
        index = faiss.read_index(os.path.join(FAISS_PATH, "faiss_index.index"))
        
        # Dense search
        query_embedding = embeddings.embed_query(query_text)
        query_embedding = np.array([query_embedding], dtype='float32')
        D_dense, I_dense = index.search(query_embedding, k=5)
        
        # Sparse search (if BM25 is available)
        if bm25:
            bm25_scores = bm25.get_scores(processed_query)
            I_sparse = np.argsort(bm25_scores)[::-1][:5]
            
            # Combine results with weighted scoring
            combined_scores = {}
            for idx, score in zip(I_dense[0], D_dense[0]):
                combined_scores[idx] = HYBRID_ALPHA * (1 - score)  # Convert distance to similarity
            
            for idx in I_sparse:
                sparse_score = bm25_scores[idx]
                combined_scores[idx] = combined_scores.get(idx, 0) + (1 - HYBRID_ALPHA) * sparse_score
            
            # Get top results
            top_indices = sorted(combined_scores.keys(), key=lambda x: combined_scores[x], reverse=True)[:5]
        else:
            top_indices = I_dense[0]
        
        # Process results
        relevant_chunks = []
        chunk_sources = {}
        
        for idx in top_indices:
            if 0 <= idx < len(chunks):
                chunk = chunks[int(idx)]
                score = combined_scores.get(idx, 0) if bm25 else D_dense[0][list(I_dense[0]).index(idx)]
                
                chunk_content = chunk.page_content
                relevant_chunks.append(chunk_content)
                
                source = chunk.metadata.get('source', 'Unknown source')
                page = chunk.metadata.get('page', 'Unknown location')
                source_name = os.path.basename(source)
                
                # Calculate relevance label based on score
                relevance = "High" if score < 0.5 else "Medium" if score < 1.0 else "Low"
                
                # Get surrounding context
                start_idx = max(0, chunk.metadata.get('start_idx', 0))
                end_idx = chunk.metadata.get('end_idx', len(chunk_content))
                
                # Store source with detailed information
                if source_name not in chunk_sources:
                    chunk_sources[source_name] = {
                        'score': score,
                        'location': f"Page {page}" if page != 'Unknown location' else 'Document',
                        'content': chunk_content[:200] + "..." if len(chunk_content) > 200 else chunk_content,
                        'relevance': relevance,
                        'context': {
                            'start': start_idx,
                            'end': end_idx
                        }
                    }
                elif score < chunk_sources[source_name]['score']:
                    # Update if we found a more relevant chunk from the same source
                    chunk_sources[source_name].update({
                        'score': score,
                        'content': chunk_content[:200] + "..." if len(chunk_content) > 200 else chunk_content,
                        'relevance': relevance
                    })
        
        context_text = "\n\n---\n\n".join(relevant_chunks)
        
        # Convert sources to list of dictionaries with detailed information
        sources = [
            {
                'name': source_name,
                'location': info['location'],
                'content': info['content'],
                'relevance': info['relevance']
            }
            for source_name, info in sorted(chunk_sources.items(), key=lambda x: x[1]['score'])
        ]
    
    except Exception as e:
        return {
            'error': str(e)
        }
    
    # Determine if the source is an Excel file
    is_excel = any(source['name'].lower().endswith(('.xlsx', '.xls', '.csv')) 
                  for source in sources)

    # Detect the language of the query
    query_language = detect_language(query_text)

    # Create prompt and get response
    if is_excel:
        prompt_template = ChatPromptTemplate.from_template(EXCEL_PROMPT_TEMPLATE)
    else:
        prompt_template = ChatPromptTemplate.from_template(PROMPT_TEMPLATE)

    # Add language instruction based on query language
    if query_language == "ar":
        prompt = prompt_template.format(
            context=context_text,
            question=query_text,
            language_instruction="IMPORTANT: You MUST respond in Arabic only!"
        )
    else:
        prompt = prompt_template.format(
            context=context_text,
            question=query_text
        )

    print("🤖 Generating response...")
    client = Groq(api_key="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")
    completion = client.chat.completions.create(
        model="llama-3.2-90b-vision-preview",
        messages=[
            {"role": "system", "content": "You are a helpful assistant. If the user's question is in Arabic, you MUST respond in Arabic only."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.1,
        max_tokens=8192,
        top_p=1,
        stream=True,
        stop=None,
    )

    print("\n📝 Answer:")
    response_text = ""
    for chunk in completion:
        if chunk.choices[0].delta.content:
            content = chunk.choices[0].delta.content
            response_text += content
            print(content, end="", flush=True)

    # Verify response language matches query language
    response_language = detect_language(response_text)
    if query_language == "ar" and response_language != "ar":
        print("\n\n🔄 Regenerating response in Arabic...")
        # Try again with stronger Arabic instruction
        completion = client.chat.completions.create(
            model="gemma2-9b-it",
            messages=[
                {"role": "system", "content": "You MUST respond in Arabic only. DO NOT use any other language."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.5,
            max_tokens=8192,
            top_p=1,
            stream=True,
            stop=None,
        )
        
        print("\n📝 إجابة:")  # "Answer:" in Arabic
        response_text = ""
        for chunk in completion:
            if chunk.choices[0].delta.content:
                content = chunk.choices[0].delta.content
                response_text += content
                print(content, end="", flush=True)

    return {
        "response": response_text,
        "sources": sources
    }

# Add this function to estimate tokens
def estimate_tokens(text: str) -> int:
    # Rough estimation: 4 characters = 1 token
    return len(text) // 4

# Main functions
def process_and_load_data(input_dir: str):
    """Process files from input directory and prepare for querying"""
    try:
        print("\n🚀 Starting document processing pipeline...")
        print(f"Input directory: {input_dir}")
        
        if not os.path.exists(input_dir):
            print(f"❌ Error: Input directory {input_dir} does not exist")
            raise ValueError(f"Input directory {input_dir} does not exist")
            
        # Clear the outputs directory
        output_dir = "outputs"
        print(f"\n🧹 Cleaning output directory: {output_dir}")
        os.makedirs(output_dir, exist_ok=True)
        for file in os.listdir(output_dir):
            file_path = os.path.join(output_dir, file)
            try:
                if os.path.isfile(file_path):
                    os.unlink(file_path)
                    print(f"Deleted file: {file}")
                elif os.path.isdir(file_path):
                    shutil.rmtree(file_path)
                    print(f"Deleted directory: {file}")
            except Exception as e:
                print(f"❌ Error cleaning output directory: {str(e)}")
        
        # Convert files to markdown
        print(f"\n📝 Converting files from {input_dir} to markdown...")
        print(f"Files in input directory: {os.listdir(input_dir)}")
        convert_files_to_md(input_dir, output_dir)
        
        # Load documents and combine their text
        print("\n📚 Loading converted markdown documents...")
        documents = load_documents()
        if not documents:
            print("❌ Error: No documents were loaded")
            raise ValueError("No documents were loaded. Please check the input files.")
            
        combined_text = "\n\n".join([doc.page_content for doc in documents])
        estimated_tokens = estimate_tokens(combined_text)
        
        print(f"\n📊 Document Statistics:")
        print(f"- Number of documents: {len(documents)}")
        print(f"- Total text length: {len(combined_text)} characters")
        print(f"- Estimated tokens: {estimated_tokens}")
        
        # Split documents into chunks
        print("\n✂️ Splitting documents into chunks...")
        chunks = split_documents(documents)
        if not chunks:
            print("❌ Error: Failed to split documents into chunks")
            raise ValueError("Failed to split documents into chunks.")
            
        # Calculate chunk IDs
        print("\n🔢 Calculating chunk IDs...")
        chunk_ids = calculate_chunk_ids(chunks)
        print(f"Generated {len(chunk_ids)} chunk IDs")
        
        # Save chunks for later use
        print("\n💾 Saving chunks...")
        os.makedirs(FAISS_PATH, exist_ok=True)
        with open(CHUNKS_PATH, 'wb') as f:
            pickle.dump({'chunks': chunks, 'chunk_ids': chunk_ids}, f)
        print(f"Chunks saved to: {CHUNKS_PATH}")
        
        # Add to FAISS
        print("\n🔍 Adding to FAISS index...")
        add_to_faiss(chunks)
        
        print("\n✅ Document processing complete!")
        return True
        
    except Exception as e:
        print(f"\n❌ Error in document processing: {str(e)}")
        traceback.print_exc()
        raise

def chat_with_documents():
    """Interactive chat loop with the documents."""
    print("\n💬 Chat with your documents! (Type 'exit' to end the conversation)")
    
    while True:
        query = input("\n❓ You: ").strip()
        
        if query.lower() in ['exit', 'quit', 'bye']:
            print("👋 Goodbye!")
            break
            
        if not query:
            continue
            
        try:
            result = query_faiss(query)
            if 'error' in result:
                print(f"❌ Error: {result['error']}")
            else:
                print("\n📝 Answer:")
                print(result['response'])
                print("\n📊 Sources:")
                for source in result['sources']:
                    print(f"- {source['name']}: {source['location']}")
        except Exception as e:
            print(f"❌ Error processing query: {str(e)}")
            print("Please try again with a different question.")

def main():
    print("🌟 Welcome to the Document Query System! 🌟")
    
    # Get input directory from user
    input_dir = input("\n📁 Enter the path to your documents folder: ").strip()

    # Process and load data from input directory
    process_and_load_data(input_dir)
    
    # Start the chat interface
    chat_with_documents()

if __name__ == "__main__":
    main()
